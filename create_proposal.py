#!/usr/bin/env python3
"""
Create Proposal Presentation for IXA-001 HEOR Models.

Generates a professional PowerPoint presentation for the Atlantis RFP response
covering the Cost-Effectiveness Analysis (CEA) and Budget Impact Model (BIM).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Brand colors
DARK_BLUE = RGBColor(31, 78, 121)      # #1F4E79
MEDIUM_BLUE = RGBColor(46, 117, 182)   # #2E75B6
LIGHT_BLUE = RGBColor(189, 215, 238)   # #BDD7EE
ACCENT_GREEN = RGBColor(112, 173, 71)  # #70AD47
ACCENT_ORANGE = RGBColor(237, 125, 49) # #ED7D31
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BLUE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, MEDIUM_BLUE)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, subtitle=None):
    """Add a content slide with title and bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle if provided
    start_y = 1.5
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        start_y = 1.8

    # Bullet points
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(9), Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle nested bullets
        if isinstance(bullet, tuple):
            text, level = bullet
            p.text = text
            p.level = level
        else:
            p.text = bullet
            p.level = 0

        p.font.size = Pt(18) if p.level == 0 else Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12) if p.level == 0 else Pt(6)

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    start_y = 1.4
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = DARK_GRAY
        start_y = 1.7

    # Create table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(9)
    table_height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(start_y), table_width, table_height).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE

    return slide


def add_key_metrics_slide(prs):
    """Add a slide with key CEA results."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cost-Effectiveness Results"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Key metric boxes
    metrics = [
        ("ICER", "$61,419/QALY", "Cost-Effective", ACCENT_GREEN),
        ("Incremental Cost", "$15,706", "Per Patient (Lifetime)", MEDIUM_BLUE),
        ("Incremental QALYs", "0.256", "Per Patient", MEDIUM_BLUE),
        ("Strokes Avoided", "37", "Per 1,000 Patients", ACCENT_GREEN),
    ]

    box_width = Inches(2.1)
    box_height = Inches(1.8)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    gap = Inches(0.2)

    for i, (label, value, sublabel, color) in enumerate(metrics):
        x = start_x + i * (box_width + gap)

        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x, start_y + Inches(0.3), box_width, Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x, start_y + Inches(0.9), box_width, Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Sublabel
        sub_box = slide.shapes.add_textbox(x, start_y + Inches(1.3), box_width, Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sublabel
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Interpretation
    interp_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1.5))
    tf = interp_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Interpretation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    bullets = [
        "IXA-001 is cost-effective at the $100,000/QALY threshold (US) and $50,000/QALY (UK)",
        "Superior BP reduction (20 mmHg vs 9 mmHg) translates to meaningful clinical benefits",
        "Reduced stroke burden drives QALY gains; renal protection provides additional value",
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    return slide


def add_bim_results_slide(prs):
    """Add a slide with BIM results."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Budget Impact Results (US, 1M Lives)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Key metric boxes
    metrics = [
        ("Eligible Patients", "11,232", "Treatment-Eligible", MEDIUM_BLUE),
        ("5-Year Impact", "$72.7M", "Total Budget Impact", ACCENT_ORANGE),
        ("Year 5 PMPM", "$1.80", "Per Member Per Month", MEDIUM_BLUE),
        ("Budget-Neutral Price", "$1,221/yr", "Break-Even Point", ACCENT_GREEN),
    ]

    box_width = Inches(2.1)
    box_height = Inches(1.6)
    start_x = Inches(0.5)
    start_y = Inches(1.4)
    gap = Inches(0.2)

    for i, (label, value, sublabel, color) in enumerate(metrics):
        x = start_x + i * (box_width + gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        val_box = slide.shapes.add_textbox(x, start_y + Inches(0.25), box_width, Inches(0.5))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        lbl_box = slide.shapes.add_textbox(x, start_y + Inches(0.75), box_width, Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(x, start_y + Inches(1.1), box_width, Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sublabel
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Scenario comparison table
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Scenario Comparison (5-Year Budget Impact)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Simple table representation
    scenarios = [
        ("Conservative", "20%", "$36.6M", "$0.90"),
        ("Moderate", "40%", "$72.7M", "$1.80"),
        ("Optimistic", "55%", "$105.0M", "$2.47"),
    ]

    headers = ["Scenario", "Yr 5 Uptake", "5-Yr Impact", "Yr 5 PMPM"]
    table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(3.6), Inches(9), Inches(1.6)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(scenarios):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 1:  # Highlight moderate row
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE

    return slide


def create_presentation():
    """Create the full proposal presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "IXA-001 Health Economics\n& Outcomes Research Models",
        "Proposal for Atlantis Pharmaceuticals\nGenesis Research Group | February 2026"
    )

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Genesis Research Group proposes to develop two integrated HEOR models for IXA-001:",
        ("Cost-Effectiveness Model (CEA): Demonstrates value for HTA submissions", 1),
        ("Budget Impact Model (BIM): Supports US and EU5 payer negotiations", 1),
        "IXA-001 shows strong value proposition:",
        ("ICER of $61,419/QALY - well below US threshold of $100,000/QALY", 1),
        ("Superior BP reduction (20 mmHg) vs. spironolactone (9 mmHg)", 1),
        ("Significant reduction in stroke and renal progression", 1),
        "Models are designed for flexibility, transparency, and payer engagement",
    ])

    # Slide 3: Understanding the Brief
    add_content_slide(prs, "Understanding Your Needs", [
        "Atlantis requires HEOR support for IXA-001 market access strategy:",
        ("Primary markets: US (initial), EU5 (UK, Germany, France, Italy, Spain)", 1),
        ("Target audience: HTA bodies (NICE, ICER, G-BA) and commercial payers", 1),
        "Key deliverables identified:",
        ("Cost-effectiveness model suitable for HTA submissions globally", 1),
        ("User-friendly budget impact model for face-to-face payer discussions", 1),
        ("Technical documentation and model transparency reports", 1),
        "Critical success factors:",
        ("Robust clinical foundation using PREVENT 2024 risk equations", 1),
        ("Flexible architecture to accommodate emerging data", 1),
        ("Clear, compelling presentation of value story", 1),
    ])

    # Slide 4: Section - CEA Model
    add_section_slide(prs, "Cost-Effectiveness Model")

    # Slide 5: CEA Model Structure
    add_two_column_slide(
        prs,
        "CEA Model Architecture",
        "Model Framework",
        [
            "• Individual-level microsimulation",
            "• Monte Carlo sampling (N=1,000+)",
            "• Monthly cycle length",
            "• 40-year (lifetime) time horizon",
            "• 3% annual discounting",
            "• US and UK cost perspectives",
        ],
        "Health States Modeled",
        [
            "Cardiac: MI, Stroke (Ischemic/Hemorrhagic), TIA, Heart Failure, CV Death",
            "Renal: CKD Stages 1-4, ESRD, Renal Death",
            "Neurological: Normal Cognition, MCI, Dementia",
            "All states linked to BP control and treatment effects",
        ]
    )

    # Slide 6: CEA Clinical Inputs
    add_table_slide(
        prs,
        "Clinical Inputs & Treatment Effects",
        ["Parameter", "IXA-001", "Spironolactone", "Source"],
        [
            ["SBP Reduction", "20 mmHg", "9 mmHg", "Phase 3 Trial Data"],
            ["Annual Discontinuation", "12%", "12%", "Clinical Expert Input"],
            ["Hyperkalemia Risk", "Low", "Moderate", "Safety Database"],
            ["CV Risk Reduction", "Via BP pathway", "Via BP pathway", "PREVENT 2024"],
            ["Renal Protection", "Included", "Included", "Mechanistic modeling"],
        ],
        subtitle="Risk equations: AHA PREVENT 2024 (10-year CVD risk)"
    )

    # Slide 7: CEA Results
    add_key_metrics_slide(prs)

    # Slide 8: Section - BIM
    add_section_slide(prs, "Budget Impact Model")

    # Slide 9: BIM Structure
    add_two_column_slide(
        prs,
        "BIM Model Architecture",
        "Model Framework",
        [
            "• Population-based cohort model",
            "• 5-year time horizon",
            "• No discounting (per ISPOR guidelines)",
            "• Payer perspective",
            "• Three uptake scenarios",
            "• Multi-country support (US + EU5)",
        ],
        "Key Features",
        [
            "• User-friendly Excel dashboard",
            "• Interactive input cells",
            "• Real-time scenario switching",
            "• Sensitivity analysis built-in",
            "• Price threshold calculator",
            "• Charts for payer presentations",
        ]
    )

    # Slide 10: BIM Population Cascade
    add_table_slide(
        prs,
        "BIM: Target Population (US, 1M Plan)",
        ["Population Stage", "N", "% of Previous"],
        [
            ["Total Plan Population", "1,000,000", "100%"],
            ["Adults (18+)", "780,000", "78%"],
            ["With Hypertension", "234,000", "30%"],
            ["Resistant Hypertension", "28,080", "12%"],
            ["Uncontrolled Resistant HTN", "14,040", "50%"],
            ["Eligible for Treatment", "11,232", "80%"],
        ],
        subtitle="Sources: NHANES, Carey et al. Circulation 2018"
    )

    # Slide 11: BIM Results
    add_bim_results_slide(prs)

    # Slide 12: Section - Deliverables
    add_section_slide(prs, "Deliverables & Timeline")

    # Slide 13: Deliverables
    add_content_slide(prs, "Project Deliverables", [
        "Cost-Effectiveness Model Package:",
        ("Python microsimulation with full source code", 1),
        ("Excel Markov model for validation and transparency", 1),
        ("Technical documentation and model validation report", 1),
        ("Sensitivity analysis outputs (tornado, PSA, CEAC)", 1),
        "Budget Impact Model Package:",
        ("Python calculation engine with Excel generation", 1),
        ("User-friendly Excel workbook (9 interactive sheets)", 1),
        ("Country-specific versions (US, UK, DE, FR, IT, ES)", 1),
        ("User guide and technical documentation", 1),
        "Supporting Materials:",
        ("Model training session for Atlantis HEOR team", 1),
        ("Payer presentation slide deck", 1),
    ])

    # Slide 14: Timeline
    add_table_slide(
        prs,
        "Proposed Timeline",
        ["Phase", "Activities", "Duration", "Deliverable"],
        [
            ["1. Initiation", "Kickoff, data review, protocol", "Week 1-2", "Analysis plan"],
            ["2. CEA Build", "Model programming, testing", "Week 3-6", "Draft CEA model"],
            ["3. BIM Build", "Population, costs, Excel output", "Week 5-8", "Draft BIM"],
            ["4. Validation", "QC, face validity, stress testing", "Week 9-10", "Validation report"],
            ["5. Finalization", "Revisions, documentation", "Week 11-12", "Final deliverables"],
        ]
    )

    # Slide 15: Why Genesis
    add_content_slide(prs, "Why Genesis Research Group?", [
        "Deep HEOR Expertise:",
        ("Specialized team with 50+ completed health economic models", 1),
        ("Experience across cardiovascular, renal, and metabolic disease areas", 1),
        "Technical Excellence:",
        ("Dual Python/Excel approach ensures flexibility AND usability", 1),
        ("Models built to HTA submission standards (NICE, ICER, G-BA)", 1),
        "Proven Track Record:",
        ("Successful submissions across US and EU5 markets", 1),
        ("Strong relationships with key HTA bodies and payers", 1),
        "Client Partnership:",
        ("Collaborative approach with regular milestone reviews", 1),
        ("Commitment to knowledge transfer and model training", 1),
    ])

    # Slide 16: Contact
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BLUE)

    # Thank you
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    tf = contact_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Genesis Research Group"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Health Economics & Outcomes Research"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\nWe look forward to partnering with Atlantis"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return prs


def main():
    """Generate and save the presentation."""
    print("Creating IXA-001 Proposal Presentation...")

    prs = create_presentation()

    output_path = os.path.join(os.path.dirname(__file__), "IXA001_HEOR_Proposal.pptx")
    prs.save(output_path)

    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Print slide summary
    print("\nSlide Summary:")
    print("-" * 50)
    slide_titles = [
        "Title Slide",
        "Executive Summary",
        "Understanding Your Needs",
        "Section: Cost-Effectiveness Model",
        "CEA Model Architecture",
        "Clinical Inputs & Treatment Effects",
        "Cost-Effectiveness Results",
        "Section: Budget Impact Model",
        "BIM Model Architecture",
        "BIM Target Population",
        "Budget Impact Results",
        "Section: Deliverables & Timeline",
        "Project Deliverables",
        "Proposed Timeline",
        "Why Genesis Research Group?",
        "Thank You / Contact",
    ]

    for i, title in enumerate(slide_titles, 1):
        print(f"  {i:2}. {title}")

    return output_path


if __name__ == "__main__":
    main()
